from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation():
    # Load original
    try:
        prs = Presentation('Mid sem Presentation (1) (1).pptx')
    except:
        prs = Presentation()
    
    # Keep only the first two slides if they exist
    # (In python-pptx, deleting slides requires accessing the XML element)
    slides = list(prs.slides)
    if len(slides) > 2:
        for i in range(2, len(slides)):
            # Access the internal slide list and remove the element
            xml_slides = prs.slides._sldIdLst
            xml_slides.remove(xml_slides[2])

    # Content for slides 3-15
    slides_content = [
        {
            "title": "The Engineering Approach",
            "bullets": [
                "Instead of just building a dashboard that displays flat data, our goal was to process information mathematically before the user even sees it.",
                "The user interface is really just the outer skin. The actual heavy lifting happens in the background algorithmic engines.",
                "We broke the problem down into distinct technical challenges: filtering bad data, grouping threats, routing units safely, and sorting alerts instantly."
            ]
        },
        {
            "title": "System Architecture & Tech Stack",
            "bullets": [
                "We built the frontend using React and TypeScript. Strict typing was absolutely necessary here because of how many raw numbers the algorithms pass around.",
                "Vite handles the build process, which kept our development cycles fast.",
                "Data currently flows from a custom-built simulation engine directly into our React hooks, mimicking a real hardware telemetry feed.",
                "The architecture is strictly decoupled—the math logic doesn't care about the UI rendering."
            ]
        },
        {
            "title": "UI/UX & Cognitive Load Management",
            "bullets": [
                "We didn't just want it to look cool; the dark 'War Room' aesthetic actually serves a functional purpose.",
                "Bright screens cause severe eye fatigue during long monitoring sessions. We stuck to a high-contrast dark mode with specific neon accents to draw the eye only when absolutely necessary.",
                "Defense Strategy: Color-coded thresholds (Red/Amber/Green) process in the human brain almost instantly, cutting down reaction time compared to complex bar charts."
            ]
        },
        {
            "title": "Core Algorithm I - Kalman Filter Mechanics",
            "bullets": [
                "The Real-World Issue: Hardware isn't perfect. If you pull raw GPS coordinates from a moving unit in the field, the signal bounces around wildly due to atmospheric interference or active electronic jamming.",
                "Why not a moving average? A simple average is entirely backward-looking. By the time a moving average smooths out a trajectory, the unit has already moved somewhere else, causing a massive visual delay on the tactical map.",
                "The Kalman Solution: We implemented a 2D Kalman Filter because it doesn't just average past data—it actually predicts the future state based on velocity vectors.",
                "The Core Equation: The entire system balances on the Kalman Gain calculation: K = P / (P + R).",
                "The Trust Ratio: Think of the gain (K) as an automated trust dial. If the incoming sensor data has too much variance or noise (R), the system drops the gain and relies on its own internal physics prediction (P).",
                "The Operational Impact: This mathematical filtering guarantees that the command dashboard shows a perfectly smooth, realistic flight path or ground track, entirely ignoring wild, corrupted data spikes."
            ]
        },
        {
            "title": "Core Algorithm II - A* Pathfinding & Risk Weights",
            "bullets": [
                "The Immediate Problem: Standard navigation tools (like Dijkstra's) are blind to danger. They will calculate a straight line directly through an active firefight just because it's technically the shortest physical distance.",
                "Why A*? We needed an algorithm that was incredibly fast for real-time movement, but flexible enough to understand terrain constraints without lagging the browser.",
                "The Math Breakdown: The core heuristic function we wrote is f(n) = g(n) + h(n) + Σ Risk(z).",
                "The Efficiency Driver (h(n)): This is our 'guess' to the finish line. It forces the search engine to pull directly toward the objective instead of flooding the map in every direction, which saves massive amounts of CPU overhead.",
                "Injecting the Risk Matrix: This is where we altered the standard algorithm. Instead of just treating a hazardous combat sector as a blocked wall, we assign it an incredibly high numerical weight (Σ Risk).",
                "The Tactical Result: Because the math is constantly trying to find the 'cheapest' path, treating danger as a high-cost variable forces the system to automatically plot safe corridors around threats, even if the physical distance is three times as long.",
                "Defense Strategy: If the panel asks why we didn't just use a simpler 'avoidance' system, point out that simple avoidance creates infinite loops where units get stuck trying to go around obstacles. The f(n) formula guarantees that the unit will eventually reach the target without getting trapped in a corner."
            ]
        },
        {
            "title": "Core Algorithm III - Sectoring via K-Means++",
            "bullets": [
                "The Issue: The battlefield gets chaotic, and command needs a way to group scattered threats into logical combat sectors.",
                "Defense Strategy: We use the '++' initialization variant to force the initial center points to be as far apart as possible, guaranteeing clean, distinct sector boundaries."
            ]
        },
        {
            "title": "Core Algorithm IV - Event Triage using Min-Heap",
            "bullets": [
                "The Issue: When fifty things happen at once, the dashboard can't freeze, and the commander can't be reading warnings out of order.",
                "Defense Strategy: Binary Min-Heap allows us to pull the most critical alert in strict O(log n) time, preventing browser thread lock-up during high-intensity alerts."
            ]
        },
        {
            "title": "Multi-Attribute Utility Theory (MAUT)",
            "bullets": [
                "We take Health, Fuel, and Ammunition and average them into a single Readiness Score.",
                "Defense Strategy: If a tank has 100% armor but 0% fuel, it is just a bunker. Our model ensures that any critical resource drop affects the whole unit's readiness state.",
                "Strict thresholds at 33% and 66% trigger the red and amber warning states."
            ]
        },
        {
            "title": "Bayesian Success Prediction",
            "bullets": [
                "The Math: P(S) = (R_avg * 0.75) - (V_avg * 0.25).",
                "Defense Strategy: This is a weighted Bayesian probability with a 3-to-1 bias to our own force readiness. A well-prepared force can overcome most environmental threats."
            ]
        },
        {
            "title": "Geospatial Integration Challenges",
            "bullets": [
                "We use Leaflet for the 2D tactical view (lightweight SVG paths) and CesiumJS for the 3D strategic perspective.",
                "Challenge: Custom synchronization logic ensures that selecting a unit in 2D tracks the exact same coordinate on the 3D globe seamlessly."
            ]
        },
        {
            "title": "Backend Mocking & State Management",
            "bullets": [
                "Simulation Engine: Maintains persistent state (fuel drain, health loss based on proximity) to test algorithms without real satellite access.",
                "Performance: We decoupled the engine from the UI, polling for updates on an interval to maintain a smooth 60fps even during heavy math processing."
            ]
        },
        {
            "title": "Edge Cases & System Testing",
            "bullets": [
                "Data Feed Drops: The Kalman filter coasts on the last known trajectory vector until connection is restored.",
                "No Safe Path: A* risk penalty maxes out and issues a 'No Viable Route' alert rather than blindly sending units into danger.",
                "Graceful Failure: The math is designed to prioritize safety over movement."
            ]
        },
        {
            "title": "Final Thoughts & Future Roadmap",
            "bullets": [
                "Proof of Concept: React/JS are fully capable of heavy algorithmic processing if memory and loops are optimized.",
                "Next Steps: Replacing simulation with actual WebSockets and upgrading Bayesian weights to a Neural Network trained on historical data."
            ]
        }
    ]

    # Add slides
    slide_layout = prs.slide_layouts[1] # Title and Content
    for content in slides_content:
        slide = prs.slides.add_slide(slide_layout)
        # Find the title and content placeholders
        title_shape = None
        content_box = None
        for shape in slide.placeholders:
            if shape.placeholder_format.type == 1: # Title
                title_shape = shape
            elif shape.placeholder_format.type == 2: # Body/Content
                content_box = shape
        
        if title_shape:
            title_shape.text = content["title"]
        else:
            slide.shapes.title.text = content["title"]
            
        if not content_box:
            # Fallback to adding a textbox if no body placeholder is found
            content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
        
        tf = content_box.text_frame
        tf.word_wrap = True
        
        for bullet in content["bullets"]:
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 0
            p.space_after = Pt(10)

    prs.save('BAS_Mid_Sem_Presentation.pptx')
    print("Success: BAS_Mid_Sem_Presentation.pptx created.")

if __name__ == "__main__":
    create_presentation()
